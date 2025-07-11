from pptx import Presentation
import comtypes.client
import os
import re
import jdatetime
import uuid
import qrcode




def generate_qr_code(data: str, save_path: str):
    qr = qrcode.QRCode(
        version=1,
        error_correction=qrcode.ERROR_CORRECT_H,
        box_size=10,
        border=4,
    )
    qr.add_data(data)
    qr.make(fit=True)
    img = qr.make_image(fill_color="#DCDCDC", back_color="white")
    with open(save_path, "wb") as f:
        img.save(f)
def add_qr_to_placeholder(slide, qr_Data, qr_output_dir="tmp"):
    """Replace QR placeholders like {{qr:...}} with QR code images"""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
            
        text = shape.text
        # Look for image placeholders in the format {{image:key}}
        image_placeholders = re.findall(r'{{qr:(.*?)}}', text)
        
        if image_placeholders:
            for img_key in image_placeholders:
                if img_key in qr_Data:
                    # Get the image path from the data
                    data = qr_Data[img_key]
                    qr_path = os.path.join(qr_output_dir, f"qr_{hash(data)}.png")
                    generate_qr_code(data, qr_path)
                    left, top, width, height = shape.left, shape.top, shape.width, shape.height
                    shape._element.getparent().remove(shape._element)
                    slide.shapes.add_picture(qr_path, left, top, width, height)
                    break  # stop after first QR in shape



def replace_placeholder_preserve_style(shape, data):
    """Replace text placeholders while preserving text formatting"""
    if not shape.has_text_frame:
        return

    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            for key, value in data.items():
                placeholder = f"{{{{{key}}}}}"
                if placeholder in run.text:
                    run.text = run.text.replace(placeholder, value)

def add_image_to_placeholder(slide, image_data):
    """Add images to placeholders with format {{image:key}}"""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
            
        text = shape.text
        # Look for image placeholders in the format {{image:key}}
        image_placeholders = re.findall(r'{{image:(.*?)}}', text)
        
        if image_placeholders:
            for img_key in image_placeholders:
                if img_key in image_data:
                    # Get the image path from the data
                    img_path = image_data[img_key]
                    
                    # Get the position and size of the placeholder shape
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height
                    
                    # Remove the shape with the placeholder
                    shape_id = shape.shape_id
                    sp = shape._element
                    sp.getparent().remove(sp)
                    
                    # Add the image at the same position and size
                    slide.shapes.add_picture(
                        img_path, 
                        left, 
                        top, 
                        width, 
                        height
                    )
                    # Once we've processed this shape, break out to avoid errors
                    # since we've removed the shape from the collection
                    break

def generate_certificate(template_path, output_dir, text_data, image_data=None,qr_data=None):
    """Generate certificate with text and image placeholders"""
    if image_data is None:
        image_data = {}
        
    prs = Presentation(template_path)
    slide = prs.slides[0]  # one template slide

    # First pass: handle all images
    if image_data:
        add_image_to_placeholder(slide, image_data)
    if qr_data:
        add_qr_to_placeholder(slide,qr_data)
    # Second pass: handle all text replacements
    for shape in slide.shapes:
        replace_placeholder_preserve_style(shape, text_data)

    output_filename = f"certificate_{text_data.get('unique', 'filled').replace(' ', '_')}"
    output_pptx = os.path.join(output_dir, f"{output_filename}.pptx")
    prs.save(output_pptx)
    return output_pptx

def convert_to_png(pptx_path, output_dir):
    """Convert PPTX to PNG"""
    powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
    powerpoint.Visible = 1

    ppt = powerpoint.Presentations.Open(pptx_path)
    
    # Get the base filename without extension
    base_filename = os.path.basename(pptx_path)
    base_filename = os.path.splitext(base_filename)[0]
    
    output_path = os.path.join(output_dir, f"{base_filename}.pdf")
    ppt.SaveAs(output_path, 32)  # 17 = PNG
    ppt.Close()
    powerpoint.Quit()
    os.remove(pptx_path)
    
    return output_path

def batch_generate_certificates(template_path, output_dir, certificate_data_list):
    """Generate multiple certificates from a list of data dictionaries"""
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    results = []
    
    for data in certificate_data_list:
        text_data = {k: v for k, v in data.items() if not k.startswith('image_')}
        image_data = {k.replace('image_', ''): v for k, v in data.items() if k.startswith('image_')}
        qr_data = {k.replace('qr_', ''): v for k, v in data.items() if k.startswith('qr_')}
        
        output_pptx = generate_certificate(template_path, output_dir, text_data, image_data,qr_data)
        output_png = convert_to_png(output_pptx, output_dir)
        
        results.append({
            'unique': text_data.get('unique', 'Unknown'),
            'pptx_path': output_pptx,
            'png_path': output_png
        })
        
        print(f"Certificate generated for {text_data.get('unique', 'Unknown')}")
    
    return results

def safe_image(image_path:str)->str:
    return image_path if os.path.exists(image_path) else "n-image.png" 

# Example usage
if __name__ == "__main__":
    template_path = "CourseTemplate.pptx"
    work_dir = os.getcwd()
    output_dir = os.path.join(work_dir, "certificates")
    date = jdatetime.datetime.now().strftime("%Y/%m/%d")
    cert_id=str(uuid.uuid4())
    # List of certificates to generate
    certificates = [
        {
            "gender":"جناب آقای",
            "name": "محمد مهدی یادگاری فرد",
            "national":"2500448258",
            "course": "امنیت سایبری",
            "org":"مرکز آموزش های آزاد",
            "date": "1404/01/25",
            "time":"25",
            "issue":date,
            "unique":cert_id.upper(),
            "number":"/".join(["404","الف","102"]),
            "signatory":"دکتر صالح حسینی",
            "position":"معاون آموزشی،پژوهشی،دانشجويی و فرهنگی",
            "image_photo":safe_image("logo.png"),
            "image_logo": safe_image("signature.png"),
            "signatory2":"دکتر صادقی خرمی",
            "position2":"مسئول مرکز آموزش های آزاد",
            "image_photo2": safe_image("logo2.png"),
            "image_logo2": safe_image("signature2.png"),
            "qr_url": f"https://my.site/cert/{cert_id}"
        }
    ]
    
    results = batch_generate_certificates(template_path, output_dir, certificates)
    
    print(f"\nAll certificates generated and saved in: {output_dir}")
    for cert in results:
        print(f"- {cert['unique']}: {cert['png_path']}")
